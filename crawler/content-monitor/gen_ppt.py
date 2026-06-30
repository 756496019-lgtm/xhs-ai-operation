from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

ICON_DIR = r"C:\Users\demiliang\Desktop\ppt_icons"

def add_icon(slide, path, x, y, size):
    """插入 PNG icon"""
    import os
    full = os.path.join(ICON_DIR, path)
    if os.path.exists(full):
        slide.shapes.add_picture(full, Cm(x), Cm(y), Cm(size), Cm(size))

prs = Presentation()
prs.slide_width  = Cm(33.867)
prs.slide_height = Cm(19.05)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── 背景色 ──────────────────────────────────────────────────────
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

# ── 颜色常量 ────────────────────────────────────────────────────
C_TEAL   = RGBColor(0x5B, 0xBF, 0xB5)
C_GREEN  = RGBColor(0x81, 0xC7, 0x84)
C_GOLD   = RGBColor(0xF6, 0xC9, 0x4E)
C_ORANGE = RGBColor(0xCC, 0x78, 0x5C)
C_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
C_PINK   = RGBColor(0xE0, 0x60, 0x90)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)
C_GREY   = RGBColor(0xAA, 0xBB, 0xCC)
C_DGREY  = RGBColor(0x88, 0x99, 0xAA)

def add_rect(slide, x, y, w, h, fill, border=None, border_w=1.5, radius=True):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE = 5 but add_shape uses shape id
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_w)
    else:
        shape.line.fill.background()
    # 圆角
    if radius:
        sp = shape.element
        prstGeom = sp.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            avLst.clear()
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 30000')
    return shape

def add_textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="微软雅黑"):
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return txb

def add_circle(slide, x, y, d, fill, border=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(9, Cm(x), Cm(y), Cm(d), Cm(d))  # 9=oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color, dash=False, w=1.5):
    from pptx.util import Cm
    connector = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(w)
    if dash:
        connector.line.dash_style = 4  # dash
    return connector

def add_card_content(slide, x, y, icon_color, items):
    """在卡片内添加带色块icon的条目"""
    for i, item in enumerate(items):
        iy = y + i * 0.75
        # 小色块 icon
        icon = slide.shapes.add_shape(1, Cm(x), Cm(iy), Cm(0.35), Cm(0.35))
        icon.fill.solid()
        icon.fill.fore_color.rgb = icon_color
        icon.line.fill.background()
        sp = icon.element
        prstGeom = sp.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
        # 文字
        add_textbox(slide, x+0.5, iy-0.1, 6.5, 0.6, item, 10, C_GREY)

# ════════════════════════════════════════════════════════════
# 标题区
# ════════════════════════════════════════════════════════════
add_textbox(slide, 0, 0.4, 33.87, 0.6, "CLAUDE CODE", 11, C_TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.0, 33.87, 1.2, "Skill 工作流架构", 30, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
# 标题下横线
line = slide.shapes.add_shape(1, Cm(10), Cm(2.35), Cm(13.87), Cm(0.06))
line.fill.solid()
line.fill.fore_color.rgb = C_TEAL
line.line.fill.background()

# ════════════════════════════════════════════════════════════
# 左侧卡片：步骤01 数据获取
# ════════════════════════════════════════════════════════════
LX, LY, LW, LH = 0.6, 4.5, 8.0, 5.8
add_rect(slide, LX, LY, LW, LH, RGBColor(0x1E, 0x3A, 0x5F), C_TEAL)
# 左侧竖条
bar = slide.shapes.add_shape(1, Cm(LX), Cm(LY), Cm(0.25), Cm(LH))
bar.fill.solid(); bar.fill.fore_color.rgb = C_TEAL; bar.line.fill.background()
# 编号圆
add_circle(slide, LX+0.5, LY+0.3, 0.85, RGBColor(0x1E, 0x4A, 0x6A), C_TEAL)
add_textbox(slide, LX+0.5, LY+0.35, 0.85, 0.85, "01", 12, C_TEAL, bold=True, align=PP_ALIGN.CENTER)
# 标题
add_textbox(slide, LX+1.5, LY+0.25, 6.0, 0.55, "DATA COLLECTION", 9, C_TEAL)
add_textbox(slide, LX+1.5, LY+0.75, 6.0, 0.8, "数据获取", 18, C_WHITE, bold=True)
# 分割线
sep = slide.shapes.add_shape(1, Cm(LX+0.4), Cm(LY+1.7), Cm(LW-0.8), Cm(0.04))
sep.fill.solid(); sep.fill.fore_color.rgb = C_TEAL; sep.line.fill.background()
# 内容
add_card_content(slide, LX+0.5, LY+2.0, C_TEAL, [
    "🌐  网页爬虫 · Playwright",
    "📡  API 数据拉取",
    "👁  用户行为 · 舆情监控",
])
# 装饰文字
add_textbox(slide, LX+5.0, LY+4.0, 2.5, 1.5, "</>", 32, C_TEAL, bold=True)
# icon
add_icon(slide, "icon_crawler.png", LX+0.4, LY+2.1, 1.2)

# ════════════════════════════════════════════════════════════
# 顶中卡片：步骤02 分析匹配
# ════════════════════════════════════════════════════════════
TX, TY, TW, TH = 13.0, 2.8, 8.0, 5.0
add_rect(slide, TX, TY, TW, TH, RGBColor(0x1A, 0x3A, 0x2A), C_GREEN)
# 顶部横条
bar2 = slide.shapes.add_shape(1, Cm(TX), Cm(TY), Cm(TW), Cm(0.25))
bar2.fill.solid(); bar2.fill.fore_color.rgb = C_GREEN; bar2.line.fill.background()
add_circle(slide, TX+0.5, TY+0.4, 0.85, RGBColor(0x1A, 0x4A, 0x2A), C_GREEN)
add_textbox(slide, TX+0.5, TY+0.45, 0.85, 0.85, "02", 12, C_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, TX+1.5, TY+0.35, 6.0, 0.55, "ANALYSIS", 9, C_GREEN)
add_textbox(slide, TX+1.5, TY+0.75, 6.0, 0.8, "分析 · 匹配", 18, C_WHITE, bold=True)
sep2 = slide.shapes.add_shape(1, Cm(TX+0.4), Cm(TY+1.7), Cm(TW-0.8), Cm(0.04))
sep2.fill.solid(); sep2.fill.fore_color.rgb = C_GREEN; sep2.line.fill.background()
add_card_content(slide, TX+0.5, TY+2.0, C_GREEN, [
    "✍  内容改写 · 提炼摘要",
    "🖼  视觉分析 · 图像理解",
    "📋  脚本理解 · 需求拆解",
])
add_textbox(slide, TX+5.5, TY+3.2, 2.0, 1.5, "✦", 30, C_GREEN, bold=True)
# icon
add_icon(slide, "icon_analysis.png", TX+0.4, TY+2.1, 1.2)

# ════════════════════════════════════════════════════════════
# 右侧卡片：步骤03 内容产出
# ════════════════════════════════════════════════════════════
RX, RY, RW, RH = 25.3, 4.5, 8.0, 5.8
add_rect(slide, RX, RY, RW, RH, RGBColor(0x3A, 0x2A, 0x1A), C_GOLD)
bar3 = slide.shapes.add_shape(1, Cm(RX+RW-0.25), Cm(RY), Cm(0.25), Cm(RH))
bar3.fill.solid(); bar3.fill.fore_color.rgb = C_GOLD; bar3.line.fill.background()
add_circle(slide, RX+0.5, RY+0.3, 0.85, RGBColor(0x4A, 0x3A, 0x1A), C_GOLD)
add_textbox(slide, RX+0.5, RY+0.35, 0.85, 0.85, "03", 12, C_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, RX+1.5, RY+0.25, 6.0, 0.55, "OUTPUT", 9, C_GOLD)
add_textbox(slide, RX+1.5, RY+0.75, 6.0, 0.8, "内容产出", 18, C_WHITE, bold=True)
sep3 = slide.shapes.add_shape(1, Cm(RX+0.4), Cm(RY+1.7), Cm(RW-0.8), Cm(0.04))
sep3.fill.solid(); sep3.fill.fore_color.rgb = C_GOLD; sep3.line.fill.background()
add_card_content(slide, RX+0.5, RY+2.0, C_GOLD, [
    "🎬  视频剪辑 · 字幕生成",
    "🎨  图文制作 · 封面设计",
    "📊  周报 · Excel · 打标",
])
add_textbox(slide, RX+5.5, RY+4.0, 2.0, 1.5, "★", 32, C_GOLD, bold=True)
# icon
add_icon(slide, "icon_output.png", RX+0.4, RY+2.1, 1.2)

# ════════════════════════════════════════════════════════════
# 中央 Skill 卡片
# ════════════════════════════════════════════════════════════
CX, CY, CW, CH = 13.0, 6.5, 8.0, 6.2
# 外圈装饰圆（用大圆形边框）
outer = slide.shapes.add_shape(9, Cm(CX-0.8), Cm(CY-0.8), Cm(CW+1.6), Cm(CH+1.6))
outer.fill.background()
outer.line.color.rgb = C_TEAL
outer.line.width = Pt(0.75)
outer.line.dash_style = 4

add_rect(slide, CX, CY, CW, CH, RGBColor(0x1E, 0x2E, 0x3E), C_TEAL, border_w=2.0)

# .md 小标签
add_textbox(slide, CX+0.3, CY+0.3, 2.0, 0.5, ".md", 9, C_TEAL)
# Skill 文件小icon（用小矩形模拟）
doc = slide.shapes.add_shape(1, Cm(CX+3.1), Cm(CY+0.5), Cm(1.8), Cm(2.2))
doc.fill.solid(); doc.fill.fore_color.rgb = RGBColor(0x24, 0x35, 0x48)
doc.line.color.rgb = C_TEAL; doc.line.width = Pt(1.0)
add_textbox(slide, CX+3.15, CY+0.9, 1.7, 0.5, "SKILL", 8, C_TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide, CX+3.15, CY+1.4, 1.7, 0.5, ".md", 9, C_TEAL, bold=True, align=PP_ALIGN.CENTER)

# 主标题
add_textbox(slide, CX, CY+2.9, CW, 1.0, "Skill", 28, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, CX, CY+3.8, CW, 0.6, "CLAUDE CODE  ·  一个文件驱动全流程", 9, C_TEAL, align=PP_ALIGN.CENTER)
# 中央 skill icon
add_icon(slide, "icon_skill.png", CX+2.8, CY+0.3, 2.4)

# 三个模型标签
tag_y = CY+4.6
for tx, label, col in [
    (CX+0.4, "Claude",  C_ORANGE),
    (CX+2.9, "Qwen",    C_PURPLE),
    (CX+5.1, "MiniMax", C_PINK),
]:
    tag = slide.shapes.add_shape(1, Cm(tx), Cm(tag_y), Cm(2.1), Cm(0.65))
    tag.fill.solid(); tag.fill.fore_color.rgb = RGBColor(0x1A, 0x30, 0x50)
    tag.line.color.rgb = col; tag.line.width = Pt(1.2)
    sp = tag.element
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
    add_textbox(slide, tx+0.05, tag_y+0.05, 2.0, 0.55, label, 9, col, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 连接箭头
# ════════════════════════════════════════════════════════════
# 左卡片 → 中央
add_line(slide, 8.6, 7.4, 13.0, 9.0, C_TEAL, dash=True)
# 上卡片 → 中央
add_line(slide, 17.0, 7.8, 17.0, 6.5, C_GREEN, dash=True)
# 中央 → 右卡片
add_line(slide, 21.0, 9.0, 25.3, 7.4, C_GOLD, dash=True)

# ════════════════════════════════════════════════════════════
# 底部三个 AI 模型卡片
# ════════════════════════════════════════════════════════════
model_y = 14.0
model_configs = [
    (5.0,  "Claude",  "C", C_ORANGE, "Anthropic · 推理/代码",  "内容生成 · 代码执行 · 逻辑推理"),
    (14.0, "Qwen",    "Q", C_PURPLE, "阿里云 · 中文/多模态",   "中文理解 · 图像识别 · 联网搜索"),
    (23.0, "MiniMax", "M", C_PINK,   "海螺AI · 语音合成",      "语音生成 · 音色克隆 · 配音"),
]
for mx, mname, mletter, mcol, msub, mdesc in model_configs:
    add_rect(slide, mx, model_y, 6.0, 3.2, RGBColor(0x1A, 0x25, 0x35), mcol, border_w=1.5)
    # 图标圆
    add_circle(slide, mx+0.35, model_y+0.4, 0.85, RGBColor(0x1A, 0x25, 0x35), mcol)
    add_textbox(slide, mx+0.35, model_y+0.45, 0.85, 0.85, mletter, 12, mcol, bold=True, align=PP_ALIGN.CENTER)
    # 名称
    add_textbox(slide, mx+1.4, model_y+0.35, 4.2, 0.6, mname, 13, C_WHITE, bold=True)
    add_textbox(slide, mx+1.4, model_y+0.9, 4.2, 0.5, msub, 9, C_GREY)
    # 分割线
    sep = slide.shapes.add_shape(1, Cm(mx+0.3), Cm(model_y+1.6), Cm(5.4), Cm(0.04))
    sep.fill.solid(); sep.fill.fore_color.rgb = mcol; sep.line.fill.background()
    # 描述
    add_textbox(slide, mx+0.35, model_y+1.75, 5.3, 0.5, mdesc, 9, C_DGREY)
    # 连接线到Skill
    add_line(slide, mx+3.0, model_y, mx+3.0 if mx==14.0 else (17.0 if mx<14.0 else 17.0), 12.7, mcol, dash=True, w=1.0)

# ════════════════════════════════════════════════════════════
# 底部说明
# ════════════════════════════════════════════════════════════
add_textbox(slide, 0, 17.6, 33.87, 0.6,
    "一个 Skill 文件，驱动完整内容生产流水线",
    11, C_TEAL, align=PP_ALIGN.CENTER)

prs.save(r"C:\Users\demiliang\Desktop\skill_workflow.pptx")
print("✅ PPT已生成：C:\\Users\\demiliang\\Desktop\\skill_workflow.pptx")
